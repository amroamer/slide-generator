import os
import uuid

from lxml import etree
from pptx import Presentation as PptxPresentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.models.presentation import Presentation
from app.models.slide import PresentationSlide
from app.services.brand_loader import load_brand_for_presentation


def _apply_rtl_paragraph(paragraph, is_rtl: bool):
    """Apply RTL direction to a python-pptx paragraph."""
    if not is_rtl:
        return
    paragraph.alignment = PP_ALIGN.RIGHT
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("rtl", "1")
    pPr.set("algn", "r")


def _apply_rtl_run(run, is_rtl: bool):
    """Set Arabic font and language on a text run."""
    if not is_rtl:
        return
    rPr = run._r.get_or_add_rPr()
    cs = rPr.find(qn("a:cs"))
    if cs is None:
        cs = etree.SubElement(rPr, qn("a:cs"))
    cs.set("typeface", "Arial")


def _set_table_rtl(table_shape):
    """Set table-level RTL direction."""
    tbl = table_shape.table._tbl
    tblPr = tbl.find(qn("a:tblPr"))
    if tblPr is None:
        tblPr = etree.SubElement(tbl, qn("a:tblPr"), {})
    tblPr.set("rtl", "1")

OUTPUT_DIR = "/app/outputs"

# Default brand colors
DEFAULT_PRIMARY = RGBColor(0x00, 0x33, 0x8D)   # KPMG Blue
DEFAULT_ACCENT = RGBColor(0x00, 0x91, 0xDA)    # KPMG Light Blue
DEFAULT_TEXT = RGBColor(0x1F, 0x29, 0x37)       # Dark gray
DEFAULT_LIGHT = RGBColor(0xF8, 0xFA, 0xFC)      # Light bg
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def _hex_to_rgb(hex_str: str) -> RGBColor:
    hex_str = hex_str.lstrip("#")
    if len(hex_str) != 6:
        return DEFAULT_PRIMARY
    return RGBColor(int(hex_str[:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=None, alignment=PP_ALIGN.LEFT, is_rtl=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text) if text else ""
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color or DEFAULT_TEXT
    p.alignment = PP_ALIGN.RIGHT if is_rtl else alignment
    if is_rtl:
        _apply_rtl_paragraph(p, True)
        for run in p.runs:
            _apply_rtl_run(run, True)
    return txBox


def _add_bullets(slide, left, top, width, height, items, font_size=14,
                 color=None, is_rtl=False, bullet_color=None):
    """Add bullet-point list with proper bullet characters, font sizes, and indentation."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    b_color = bullet_color or DEFAULT_ACCENT

    for i, item in enumerate(items or []):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Set bullet character via XML
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(Emu(Inches(0.35))))   # left margin
        pPr.set("indent", str(Emu(-Inches(0.25)))) # hanging indent

        # Add bullet character element
        buChar = etree.SubElement(pPr, qn("a:buChar"))
        buChar.set("char", "\u2022")  # bullet dot

        # Set bullet color
        buClr = etree.SubElement(pPr, qn("a:buClr"))
        srgbClr = etree.SubElement(buClr, qn("a:srgbClr"))
        srgbClr.set("val", f"{b_color[0]:02X}{b_color[1]:02X}{b_color[2]:02X}"
                     if hasattr(b_color, '__getitem__') else "0091DA")

        # Set bullet size relative to text
        buSzPct = etree.SubElement(pPr, qn("a:buSzPct"))
        buSzPct.set("val", "100000")  # 100%

        # Add text run
        run = p.add_run()
        run.text = str(item) if item else ""
        run.font.size = Pt(font_size)
        run.font.color.rgb = color or DEFAULT_TEXT

        p.space_after = Pt(6)
        p.level = 0

        if is_rtl:
            _apply_rtl_paragraph(p, True)
            _apply_rtl_run(run, True)

    return txBox


def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def build_title_slide(pptx_slide, content, primary, accent, is_rtl=False):
    _set_slide_bg(pptx_slide, primary)
    title = content.get("title", "")
    align = PP_ALIGN.CENTER
    _add_textbox(pptx_slide, Inches(1.5), Inches(2.5), Inches(10), Inches(2),
                 title, font_size=36, bold=True, color=WHITE, alignment=align, is_rtl=is_rtl)
    subtitle = (content.get("body", {}) or {}).get("content", [""])[0] if content.get("body") else ""
    if subtitle:
        _add_textbox(pptx_slide, Inches(2), Inches(4.5), Inches(9), Inches(1),
                     subtitle, font_size=18, color=WHITE, alignment=align, is_rtl=is_rtl)
    _add_shape_bar(pptx_slide, Inches(4), Inches(4.2), Inches(5), Inches(0.05), accent)


def build_title_bullets_slide(pptx_slide, content, primary, accent, is_rtl=False):
    _add_shape_bar(pptx_slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), primary)
    _add_textbox(pptx_slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 content.get("title", ""), font_size=28, bold=True, color=primary, is_rtl=is_rtl)
    bar_left = Inches(10.5) if is_rtl else Inches(0.8)
    _add_shape_bar(pptx_slide, bar_left, Inches(1.15), Inches(2), Inches(0.04), accent)
    body = (content.get("body", {}) or {}).get("content", [])
    _add_bullets(pptx_slide, Inches(0.8), Inches(1.5), Inches(11), Inches(5),
                 body, font_size=16, is_rtl=is_rtl, bullet_color=accent)
    kt = content.get("key_takeaway")
    if kt:
        kt_bar_left = Inches(12.4) if is_rtl else Inches(0.7)
        _add_shape_bar(pptx_slide, kt_bar_left, Inches(6.3), Inches(0.08), Inches(0.6), accent)
        _add_textbox(pptx_slide, Inches(1.0), Inches(6.3), Inches(10), Inches(0.6),
                     kt, font_size=13, bold=True, color=primary, is_rtl=is_rtl)


# RAG / semantic color maps
SEMANTIC_COLORS_MAP = {
    "green": RGBColor(0x10, 0xB9, 0x81), "on track": RGBColor(0x10, 0xB9, 0x81),
    "on-track": RGBColor(0x10, 0xB9, 0x81), "met": RGBColor(0x10, 0xB9, 0x81),
    "healthy": RGBColor(0x10, 0xB9, 0x81), "completed": RGBColor(0x10, 0xB9, 0x81),
    "pass": RGBColor(0x10, 0xB9, 0x81), "yes": RGBColor(0x10, 0xB9, 0x81),
    "positive": RGBColor(0x10, 0xB9, 0x81), "good": RGBColor(0x10, 0xB9, 0x81),
    "low": RGBColor(0x10, 0xB9, 0x81),
    "amber": RGBColor(0xF5, 0x9E, 0x0B), "yellow": RGBColor(0xF5, 0x9E, 0x0B),
    "warning": RGBColor(0xF5, 0x9E, 0x0B), "attention": RGBColor(0xF5, 0x9E, 0x0B),
    "at risk": RGBColor(0xF5, 0x9E, 0x0B), "at-risk": RGBColor(0xF5, 0x9E, 0x0B),
    "caution": RGBColor(0xF5, 0x9E, 0x0B), "in progress": RGBColor(0xF5, 0x9E, 0x0B),
    "partial": RGBColor(0xF5, 0x9E, 0x0B), "medium": RGBColor(0xF5, 0x9E, 0x0B),
    "red": RGBColor(0xEF, 0x44, 0x44), "critical": RGBColor(0xEF, 0x44, 0x44),
    "fail": RGBColor(0xEF, 0x44, 0x44), "failed": RGBColor(0xEF, 0x44, 0x44),
    "not met": RGBColor(0xEF, 0x44, 0x44), "off track": RGBColor(0xEF, 0x44, 0x44),
    "off-track": RGBColor(0xEF, 0x44, 0x44), "overdue": RGBColor(0xEF, 0x44, 0x44),
    "no": RGBColor(0xEF, 0x44, 0x44), "blocked": RGBColor(0xEF, 0x44, 0x44),
    "negative": RGBColor(0xEF, 0x44, 0x44), "high": RGBColor(0xEF, 0x44, 0x44),
    "neutral": RGBColor(0x6B, 0x72, 0x80),
}

# Light background versions for table cell fills
SEMANTIC_BG_MAP = {
    "green": RGBColor(0xD1, 0xFA, 0xE5), "on track": RGBColor(0xD1, 0xFA, 0xE5),
    "on-track": RGBColor(0xD1, 0xFA, 0xE5), "met": RGBColor(0xD1, 0xFA, 0xE5),
    "healthy": RGBColor(0xD1, 0xFA, 0xE5), "completed": RGBColor(0xD1, 0xFA, 0xE5),
    "pass": RGBColor(0xD1, 0xFA, 0xE5), "yes": RGBColor(0xD1, 0xFA, 0xE5),
    "positive": RGBColor(0xD1, 0xFA, 0xE5), "good": RGBColor(0xD1, 0xFA, 0xE5),
    "low": RGBColor(0xD1, 0xFA, 0xE5),
    "amber": RGBColor(0xFE, 0xF3, 0xC7), "yellow": RGBColor(0xFE, 0xF3, 0xC7),
    "warning": RGBColor(0xFE, 0xF3, 0xC7), "attention": RGBColor(0xFE, 0xF3, 0xC7),
    "at risk": RGBColor(0xFE, 0xF3, 0xC7), "at-risk": RGBColor(0xFE, 0xF3, 0xC7),
    "caution": RGBColor(0xFE, 0xF3, 0xC7), "in progress": RGBColor(0xFE, 0xF3, 0xC7),
    "partial": RGBColor(0xFE, 0xF3, 0xC7), "medium": RGBColor(0xFE, 0xF3, 0xC7),
    "red": RGBColor(0xFE, 0xE2, 0xE2), "critical": RGBColor(0xFE, 0xE2, 0xE2),
    "fail": RGBColor(0xFE, 0xE2, 0xE2), "failed": RGBColor(0xFE, 0xE2, 0xE2),
    "not met": RGBColor(0xFE, 0xE2, 0xE2), "off track": RGBColor(0xFE, 0xE2, 0xE2),
    "off-track": RGBColor(0xFE, 0xE2, 0xE2), "overdue": RGBColor(0xFE, 0xE2, 0xE2),
    "no": RGBColor(0xFE, 0xE2, 0xE2), "blocked": RGBColor(0xFE, 0xE2, 0xE2),
    "negative": RGBColor(0xFE, 0xE2, 0xE2), "high": RGBColor(0xFE, 0xE2, 0xE2),
    "neutral": RGBColor(0xF3, 0xF4, 0xF6),
}

# Text colors for RAG cells (dark on light background)
SEMANTIC_TEXT_MAP = {
    "green": RGBColor(0x06, 0x5F, 0x46),
    "amber": RGBColor(0x92, 0x40, 0x0E),
    "red": RGBColor(0x99, 0x1B, 0x1B),
    "neutral": RGBColor(0x37, 0x41, 0x51),
}


def _get_rag_category(value: str) -> str | None:
    """Return 'green', 'amber', 'red', or None based on cell value."""
    norm = value.lower().strip()
    for key in SEMANTIC_COLORS_MAP:
        if key == norm:
            if key in ("green", "on track", "on-track", "met", "healthy", "completed",
                       "pass", "yes", "positive", "good", "low"):
                return "green"
            elif key in ("amber", "yellow", "warning", "attention", "at risk", "at-risk",
                         "caution", "in progress", "partial", "medium"):
                return "amber"
            elif key in ("red", "critical", "fail", "failed", "not met", "off track",
                         "off-track", "overdue", "no", "blocked", "negative", "high"):
                return "red"
            else:
                return "neutral"
    return None


def _detect_semantic_colors(labels: list[str]) -> dict[str, RGBColor] | None:
    """Detect semantic meaning in chart labels and return color map."""
    color_map: dict[str, RGBColor] = {}
    match_count = 0
    for label in labels:
        norm = label.lower().strip()
        for key, color in SEMANTIC_COLORS_MAP.items():
            if key in norm:
                color_map[label] = color
                match_count += 1
                break
    return color_map if match_count > len(labels) * 0.5 else None


CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "vertical_bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "horizontal_bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "area": XL_CHART_TYPE.AREA,
    "pie": XL_CHART_TYPE.PIE,
    "donut": XL_CHART_TYPE.DOUGHNUT,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
}

SERIES_COLORS = [
    RGBColor(0x00, 0x33, 0x8D), RGBColor(0x00, 0x91, 0xDA),
    RGBColor(0x48, 0x36, 0x98), RGBColor(0x00, 0xA3, 0xA1),
    RGBColor(0xC6, 0x00, 0x7E), RGBColor(0xFF, 0x6D, 0x00),
]


def build_title_table_slide(pptx_slide, content, primary, accent, is_rtl=False):
    _add_shape_bar(pptx_slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), primary)
    _add_textbox(pptx_slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 content.get("title", ""), font_size=28, bold=True, color=primary, is_rtl=is_rtl)

    table_data = content.get("data_table")
    if table_data and table_data.get("headers"):
        headers = table_data["headers"]
        rows_data = table_data.get("rows", [])
        n_rows = len(rows_data) + 1
        n_cols = len(headers)
        table_shape = pptx_slide.shapes.add_table(
            n_rows, n_cols, Inches(0.8), Inches(1.5),
            Inches(11.5), Inches(min(5.5, 0.4 * n_rows + 0.5))
        )
        if is_rtl:
            _set_table_rtl(table_shape)
        tbl = table_shape.table

        # Style header row
        for ci, h in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.text = str(h)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = WHITE
                _apply_rtl_paragraph(p, is_rtl)
            cell.fill.solid()
            cell.fill.fore_color.rgb = primary

        # Style data rows with RAG color detection
        for ri, row in enumerate(rows_data):
            for ci, val in enumerate(row):
                if ci < n_cols:
                    cell = tbl.cell(ri + 1, ci)
                    cell.text = str(val)

                    # Check for RAG semantic color
                    rag = _get_rag_category(str(val))
                    if rag:
                        bg = SEMANTIC_BG_MAP.get(rag)
                        txt_color = SEMANTIC_TEXT_MAP.get(rag, DEFAULT_TEXT)
                        if bg:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = bg
                        for p in cell.text_frame.paragraphs:
                            p.font.size = Pt(10)
                            p.font.bold = True
                            p.font.color.rgb = txt_color
                            p.alignment = PP_ALIGN.CENTER
                    else:
                        for p in cell.text_frame.paragraphs:
                            p.font.size = Pt(10)
                        if ri % 2 == 1:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = DEFAULT_LIGHT
    else:
        body = (content.get("body", {}) or {}).get("content", [])
        _add_bullets(pptx_slide, Inches(0.8), Inches(1.5), Inches(11), Inches(5),
                     body, bullet_color=accent)


def build_title_chart_slide(pptx_slide, content, primary, accent, is_rtl=False):
    _add_shape_bar(pptx_slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), primary)
    _add_textbox(pptx_slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 content.get("title", ""), font_size=28, bold=True, color=primary, is_rtl=is_rtl)

    chart_raw = content.get("chart_data")
    if not chart_raw or not chart_raw.get("labels") or not chart_raw.get("datasets"):
        body = (content.get("body", {}) or {}).get("content", []) or []
        if body:
            _add_bullets(pptx_slide, Inches(0.8), Inches(1.5), Inches(11), Inches(5),
                         body, bullet_color=accent)
        else:
            _add_textbox(pptx_slide, Inches(3), Inches(3.5), Inches(7), Inches(1),
                         "Chart data not available", font_size=16, color=accent, alignment=PP_ALIGN.CENTER)
        return

    labels = chart_raw["labels"]
    datasets = chart_raw["datasets"]
    raw_type = (chart_raw.get("chart_type") or "bar").lower().replace(" ", "_")
    xl_type = CHART_TYPE_MAP.get(raw_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    cd = CategoryChartData()
    cd.categories = [str(l) for l in labels]
    for ds in datasets:
        values = []
        for v in (ds.get("values") or []):
            try:
                values.append(float(v))
            except (ValueError, TypeError):
                values.append(0)
        cd.add_series(str(ds.get("label", "Series")), values)

    chart_left = Inches(0.8)
    chart_top = Inches(1.5)
    chart_width = Inches(11.5)
    chart_height = Inches(5.0)

    chart_frame = pptx_slide.shapes.add_chart(
        xl_type, chart_left, chart_top, chart_width, chart_height, cd
    )
    chart = chart_frame.chart
    chart.has_legend = len(datasets) > 1

    # Detect semantic colors
    semantic = _detect_semantic_colors(labels)

    # Style series/point colors
    is_pie = raw_type in ("pie", "donut", "doughnut")
    if semantic and (is_pie or len(datasets) == 1):
        for series in chart.series:
            for pt_idx, label in enumerate(labels):
                try:
                    point = series.points[pt_idx]
                    pt_fill = point.format.fill
                    pt_fill.solid()
                    pt_fill.fore_color.rgb = semantic.get(
                        label, SERIES_COLORS[pt_idx % len(SERIES_COLORS)]
                    )
                except (IndexError, Exception):
                    pass
    else:
        for i, series in enumerate(chart.series):
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = SERIES_COLORS[i % len(SERIES_COLORS)]

    chart.has_legend = len(datasets) > 1 or is_pie

    # Style fonts
    try:
        if chart.category_axis:
            chart.category_axis.tick_labels.font.size = Pt(8)
        if chart.value_axis:
            chart.value_axis.tick_labels.font.size = Pt(8)
    except Exception:
        pass  # Pie/donut charts don't have category/value axes

    # Key takeaway below chart
    kt = content.get("key_takeaway")
    if kt:
        _add_shape_bar(pptx_slide, Inches(0.7), Inches(6.7), Inches(0.08), Inches(0.5), accent)
        _add_textbox(pptx_slide, Inches(1.0), Inches(6.7), Inches(10), Inches(0.5),
                     kt, font_size=11, bold=True, color=primary)


def build_two_column_slide(pptx_slide, content, primary, accent, is_rtl=False):
    _add_shape_bar(pptx_slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), primary)
    _add_textbox(pptx_slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 content.get("title", ""), font_size=28, bold=True, color=primary, is_rtl=is_rtl)
    body = (content.get("body", {}) or {}).get("content", []) or []
    mid = max(len(body) // 2, 1) if body else 1
    if is_rtl:
        _add_bullets(pptx_slide, Inches(7), Inches(1.5), Inches(5.5), Inches(5),
                     body[:mid], is_rtl=True, bullet_color=accent)
        _add_shape_bar(pptx_slide, Inches(6.5), Inches(1.5), Inches(0.02), Inches(5), RGBColor(0xE5, 0xE7, 0xEB))
        _add_bullets(pptx_slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5),
                     body[mid:], is_rtl=True, bullet_color=accent)
    else:
        _add_bullets(pptx_slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5),
                     body[:mid], bullet_color=accent)
        _add_shape_bar(pptx_slide, Inches(6.5), Inches(1.5), Inches(0.02), Inches(5), RGBColor(0xE5, 0xE7, 0xEB))
        _add_bullets(pptx_slide, Inches(7), Inches(1.5), Inches(5.5), Inches(5),
                     body[mid:], bullet_color=accent)


def build_section_divider_slide(pptx_slide, content, primary, accent, is_rtl=False):
    _set_slide_bg(pptx_slide, primary)
    _add_textbox(pptx_slide, Inches(1.5), Inches(2.8), Inches(10), Inches(1.5),
                 content.get("title", ""), font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER, is_rtl=is_rtl)


def build_key_takeaway_slide(pptx_slide, content, primary, accent, is_rtl=False):
    _add_shape_bar(pptx_slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), primary)
    _add_textbox(pptx_slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 content.get("title", ""), font_size=28, bold=True, color=primary, is_rtl=is_rtl)
    kt = content.get("key_takeaway", "")
    if kt:
        kt_bar_left = Inches(11.5) if is_rtl else Inches(1.5)
        _add_shape_bar(pptx_slide, kt_bar_left, Inches(2), Inches(0.1), Inches(1.5), accent)
        _add_textbox(pptx_slide, Inches(2), Inches(2), Inches(9), Inches(1.5),
                     kt, font_size=22, bold=True, color=primary, is_rtl=is_rtl)
    body = (content.get("body", {}) or {}).get("content", [])
    _add_bullets(pptx_slide, Inches(1.5), Inches(4), Inches(10), Inches(3),
                 body, is_rtl=is_rtl, bullet_color=accent)


LAYOUT_BUILDERS = {
    "title_slide": build_title_slide,
    "title_bullets": build_title_bullets_slide,
    "title_table": build_title_table_slide,
    "title_chart": build_title_chart_slide,
    "two_column": build_two_column_slide,
    "section_divider": build_section_divider_slide,
    "key_takeaway": build_key_takeaway_slide,
    "full_image": build_title_bullets_slide,  # fallback
}


def _add_brand_logo(slide, brand):
    """Add brand logo to slide based on position and size settings."""
    size_map = {"small": Inches(0.8), "medium": Inches(1.2), "large": Inches(1.6)}
    logo_w = size_map.get(brand.logo_size, Inches(1.2))
    pos = brand.logo_position or "top-right"
    margin = Inches(0.3)
    if "right" in pos:
        left = SLIDE_WIDTH - logo_w - margin
    else:
        left = margin
    if "bottom" in pos:
        top = SLIDE_HEIGHT - logo_w * 0.6 - margin
    else:
        top = margin
    try:
        slide.shapes.add_picture(brand.logo_path, left, top, width=logo_w)
    except Exception:
        pass


def _remove_slide_transition(pptx_slide):
    """Remove any slide transition by clearing the <p:transition> element."""
    cSld = pptx_slide._element
    # Remove existing transition element
    for trans in cSld.findall(qn("p:transition")):
        cSld.remove(trans)
    # Add an explicit empty transition with advClick only (no animation)
    trans = etree.SubElement(cSld, qn("p:transition"))
    trans.set("advClick", "1")
    trans.set("spd", "fast")


def _remove_shape_animations(pptx_slide):
    """Remove all shape-level animations from a slide."""
    cSld = pptx_slide._element
    # Remove <p:timing> element which contains all animations
    for timing in cSld.findall(qn("p:timing")):
        cSld.remove(timing)


async def generate_pptx(
    presentation_id: uuid.UUID,
    slide_ids: list[str] | None,
    db: AsyncSession,
) -> str:
    pres = (await db.execute(
        select(Presentation).where(Presentation.id == presentation_id)
    )).scalar_one()

    query = (
        select(PresentationSlide)
        .where(PresentationSlide.presentation_id == presentation_id)
        .order_by(PresentationSlide.order)
    )
    if slide_ids:
        query = query.where(PresentationSlide.slide_id.in_(slide_ids))
    slides = (await db.execute(query)).scalars().all()

    # Load brand profile
    brand = await load_brand_for_presentation(presentation_id, db)
    if brand:
        primary = _hex_to_rgb(brand.primary_color)
        accent = _hex_to_rgb(brand.secondary_color)
    else:
        primary = DEFAULT_PRIMARY
        accent = DEFAULT_ACCENT

    # Detect RTL from presentation language
    lang = pres.language.value if hasattr(pres.language, "value") else str(pres.language)
    is_rtl = lang in ("arabic", "bilingual")

    pptx = PptxPresentation()
    pptx.slide_width = SLIDE_WIDTH
    pptx.slide_height = SLIDE_HEIGHT

    blank_layout = pptx.slide_layouts[6]  # Blank layout

    for sl in slides:
        pptx_slide = pptx.slides.add_slide(blank_layout)
        content = sl.content_json or {}
        layout = sl.layout or "title_bullets"

        # Auto-detect: if slide has chart_data, use chart builder regardless of layout
        chart_d = content.get("chart_data")
        if chart_d and chart_d.get("labels") and chart_d.get("datasets"):
            builder = build_title_chart_slide
        elif content.get("data_table") and content["data_table"].get("headers"):
            builder = build_title_table_slide
        else:
            builder = LAYOUT_BUILDERS.get(layout, build_title_bullets_slide)

        builder(pptx_slide, content, primary, accent, is_rtl=is_rtl)

        # Add brand logo if available
        if brand and brand.logo_path and os.path.exists(brand.logo_path):
            _add_brand_logo(pptx_slide, brand)

        # Speaker notes
        notes = content.get("speaker_notes", "")
        if notes:
            notes_slide = pptx_slide.notes_slide
            notes_slide.notes_text_frame.text = notes

        # Remove transitions and animations — clean, static slides
        _remove_slide_transition(pptx_slide)
        _remove_shape_animations(pptx_slide)

    # Save
    out_dir = os.path.join(OUTPUT_DIR, str(presentation_id))
    os.makedirs(out_dir, exist_ok=True)
    safe_title = pres.title.replace("/", "_").replace("\\", "_")[:80]
    filename = f"{safe_title}.pptx"
    filepath = os.path.join(out_dir, filename)
    pptx.save(filepath)
    return filepath

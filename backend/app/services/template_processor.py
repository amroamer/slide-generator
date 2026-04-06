"""Process uploaded PPTX files into template collections with variations."""

import logging
import os
import uuid

from lxml import etree
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation as PptxPresentation
from pptx.util import Emu
from sqlalchemy.ext.asyncio import AsyncSession

from app.models.slide_template import TemplateCollection, TemplateVariation
from app.services.template_quality import compute_variation_metrics

logger = logging.getLogger(__name__)

TEMPLATE_DIR = "/app/uploads/templates"
THUMB_W, THUMB_H = 960, 540  # 16:9 at reasonable resolution


def _emu_to_inches(emu_val) -> float:
    if emu_val is None:
        return 0
    return round(emu_val / 914400, 2)


def _color_to_hex(color_obj) -> str | None:
    try:
        if color_obj and color_obj.rgb:
            return str(color_obj.rgb)
    except Exception:
        pass
    return None


def _extract_fill_color(shape) -> str | None:
    try:
        fill = shape.fill
        if fill.type is not None:
            fc = fill.fore_color
            if fc and fc.rgb:
                return str(fc.rgb)
    except Exception:
        pass
    return None


def _is_placeholder_text(text: str) -> bool:
    if not text or len(text.strip()) < 3:
        return True
    patterns = [
        "lorem ipsum", "click to add", "add text", "your text",
        "title here", "subtitle", "body text", "description",
        "insert text", "placeholder", "sample text",
        "item 1", "item 2", "point 1", "point 2",
        "milestone", "phase 1", "phase 2",
        "xx%", "xxx", "##",
    ]
    lower = text.lower().strip()
    return any(p in lower for p in patterns)


def _identify_slot_type(shape_info: dict, idx: int, total: int) -> str | None:
    text = shape_info.get("text", "").lower()
    fs = shape_info.get("font_size_pt") or 0
    bold = shape_info.get("font_bold", False)
    top = shape_info.get("position", {}).get("top_inches", 0)

    if fs >= 20 and top < 2:
        return "title"
    if bold and fs >= 16 and top < 2.5:
        return "title"
    if fs >= 12 and top < 3 and not bold:
        return "subtitle"
    if fs <= 10 and len(text) < 20:
        return "label"
    if text:
        return "item"
    return None


def _detect_layout_style(shapes: list) -> str:
    if not shapes:
        return "centered"
    positions = [(s["position"]["left_inches"], s["position"]["top_inches"]) for s in shapes if s.get("position")]
    if not positions:
        return "centered"
    lefts = [p[0] for p in positions]
    tops = [p[1] for p in positions]
    top_var = max(tops) - min(tops) if tops else 0
    left_var = max(lefts) - min(lefts) if lefts else 0

    if top_var < 2 and left_var > 5:
        return "horizontal_flow"
    if left_var < 2 and top_var > 3:
        return "vertical_flow"
    if top_var < 1.5 and left_var < 3:
        return "centered"

    unique_tops = len(set(round(t, 1) for t in tops))
    unique_lefts = len(set(round(l, 1) for l in lefts))
    if unique_tops >= 2 and unique_lefts >= 2:
        return "grid"
    return "scattered"


def _auto_name(design: dict, index: int) -> str:
    style_map = {
        "horizontal_flow": "Horizontal", "vertical_flow": "Vertical",
        "grid": "Grid", "centered": "Centered", "scattered": "Free Layout",
    }
    style = style_map.get(design.get("layout_style", ""), "Custom")
    letter = chr(65 + index)
    return f"Style {letter} \u2014 {style}"


def _auto_tag(design: dict) -> list:
    tags = []
    colors = design.get("color_palette", [])
    shapes = design.get("shapes", [])
    if len(colors) <= 2:
        tags.append("minimalist")
    elif len(colors) >= 5:
        tags.append("colorful")
    if len(shapes) <= 5:
        tags.append("clean")
    elif len(shapes) >= 15:
        tags.append("detailed")
    tags.append(design.get("layout_style", "custom"))
    return tags


def _hex_to_rgb_tuple(hex_str: str | None) -> tuple:
    if not hex_str:
        return (200, 200, 200)
    hex_str = hex_str.lstrip("#")
    if len(hex_str) != 6:
        return (200, 200, 200)
    return (int(hex_str[:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def _get_font(size: int, bold: bool = False) -> ImageFont.ImageFont:
    """Get the best available font."""
    candidates = [
        f"/usr/share/fonts/truetype/liberation/LiberationSans-{'Bold' if bold else 'Regular'}.ttf",
        f"/usr/share/fonts/truetype/dejavu/DejaVuSans{'-Bold' if bold else ''}.ttf",
        "/usr/share/fonts/truetype/noto/NotoSans-Regular.ttf",
    ]
    for path in candidates:
        try:
            return ImageFont.truetype(path, size)
        except Exception:
            continue
    return ImageFont.load_default()


def generate_thumbnails_batch(pptx_path: str, collection_id: str, slide_count: int) -> list[str]:
    """Generate thumbnails using LibreOffice PDF conversion, with Python fallback."""
    import subprocess
    import shutil
    from pathlib import Path

    output_dir = os.path.join(TEMPLATE_DIR, str(collection_id), "thumbnails")
    os.makedirs(output_dir, exist_ok=True)

    # METHOD 1: LibreOffice → PDF → PNG (best quality)
    try:
        temp_dir = f"/tmp/template_pdf_{collection_id}"
        os.makedirs(temp_dir, exist_ok=True)

        lo_home = f"/tmp/lo_home_{collection_id}"
        os.makedirs(lo_home, exist_ok=True)
        env = os.environ.copy()
        env["HOME"] = lo_home

        logger.info("Converting PPTX to PDF via LibreOffice: %s", pptx_path)
        result = subprocess.run(
            ["libreoffice", "--headless", "--norestore", "--nofirststartwizard",
             "--nologo", "--convert-to", "pdf", "--outdir", temp_dir, pptx_path],
            capture_output=True, text=True, timeout=180, env=env,
        )
        logger.info("LibreOffice stdout: %s", result.stdout)
        if result.stderr:
            logger.warning("LibreOffice stderr: %s", result.stderr)

        pdf_files = list(Path(temp_dir).glob("*.pdf"))
        if pdf_files:
            pdf_path = str(pdf_files[0])
            logger.info("PDF generated: %s (%d bytes)", pdf_path, os.path.getsize(pdf_path))

            # Convert each page to PNG
            thumbnail_paths = []
            for i in range(slide_count):
                thumb_base = os.path.join(output_dir, f"slide_{i}")
                thumb_path = f"{thumb_base}.png"
                try:
                    subprocess.run(
                        ["pdftoppm", "-png", "-r", "200", "-f", str(i + 1), "-l", str(i + 1),
                         "-singlefile", pdf_path, thumb_base],
                        capture_output=True, text=True, timeout=30,
                    )
                    if os.path.exists(thumb_path):
                        # Verify and normalize to 16:9
                        img = Image.open(thumb_path)
                        aspect = img.width / img.height if img.height > 0 else 0
                        if aspect < 1.2:
                            # Wrong aspect — retry without cropbox
                            os.remove(thumb_path)
                            subprocess.run(
                                ["pdftoppm", "-png", "-r", "200", "-f", str(i + 1), "-l", str(i + 1),
                                 "-singlefile", pdf_path, thumb_base],
                                capture_output=True, text=True, timeout=30,
                            )
                        if os.path.exists(thumb_path):
                            # Resize to standard
                            img = Image.open(thumb_path)
                            if img.width > 1920:
                                img.thumbnail((1920, 1080), Image.LANCZOS)
                                img.save(thumb_path, "PNG")
                            logger.info("Thumbnail slide_%d: %dx%d", i, img.width, img.height)
                            thumbnail_paths.append(thumb_path)
                        else:
                            thumbnail_paths.append(_generate_python_thumbnail(pptx_path, i, output_dir))
                    else:
                        # Check for alternative name
                        alts = list(Path(output_dir).glob(f"slide_{i}*.png"))
                        if alts:
                            os.rename(str(alts[0]), thumb_path)
                            thumbnail_paths.append(thumb_path)
                        else:
                            thumbnail_paths.append(_generate_python_thumbnail(pptx_path, i, output_dir))
                except Exception as e:
                    logger.error("pdftoppm failed for slide %d: %s", i, e)
                    thumbnail_paths.append(_generate_python_thumbnail(pptx_path, i, output_dir))

            # Cleanup
            shutil.rmtree(temp_dir, ignore_errors=True)
            shutil.rmtree(lo_home, ignore_errors=True)

            # Verify quality
            good = sum(1 for p in thumbnail_paths if os.path.exists(p) and os.path.getsize(p) > 5000)
            if good >= slide_count * 0.5:
                logger.info("LibreOffice thumbnails OK: %d/%d good", good, slide_count)
                return thumbnail_paths
            logger.warning("LibreOffice thumbnails low quality, falling back to Python")

    except FileNotFoundError:
        logger.info("LibreOffice not available, using Python renderer")
    except subprocess.TimeoutExpired:
        logger.error("LibreOffice conversion timed out")
    except Exception as e:
        logger.error("LibreOffice conversion failed: %s", e)

    # METHOD 2: Python-based rendering from PPTX shape data
    logger.info("Generating thumbnails with Python renderer")
    paths = []
    for i in range(slide_count):
        paths.append(_generate_python_thumbnail(pptx_path, i, output_dir))
    return paths


def _generate_python_thumbnail(pptx_path: str, slide_index: int, output_dir: str) -> str:
    """Render a slide to PNG using python-pptx shape data + Pillow."""
    thumb_path = os.path.join(output_dir, f"slide_{slide_index}.png")

    try:
        prs = PptxPresentation(pptx_path)
        slide = prs.slides[slide_index]
        slide_w = prs.slide_width or 12192000
        slide_h = prs.slide_height or 6858000
        scale_x = THUMB_W / slide_w
        scale_y = THUMB_H / slide_h

        img = Image.new("RGB", (THUMB_W, THUMB_H), (255, 255, 255))
        draw = ImageDraw.Draw(img)

        # Background
        try:
            bg = slide.background
            if bg and bg.fill and bg.fill.type is not None:
                bg_rgb = _hex_to_rgb_tuple(str(bg.fill.fore_color.rgb))
                draw.rectangle([0, 0, THUMB_W, THUMB_H], fill=bg_rgb)
        except Exception:
            pass

        # Render shapes
        for shape in slide.shapes:
            try:
                _render_shape_to_image(draw, img, shape, scale_x, scale_y)
            except Exception:
                continue

        draw.rectangle([0, 0, THUMB_W - 1, THUMB_H - 1], outline=(230, 230, 230), width=1)
        img.save(thumb_path, "PNG", quality=95)
        return thumb_path

    except Exception as e:
        logger.error("Python thumbnail failed for slide %d: %s", slide_index, e)
        return _generate_placeholder(output_dir, slide_index)


def _render_shape_to_image(draw, img, shape, sx, sy):
    """Render a single PPTX shape onto the PIL image."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from io import BytesIO

    x = int((shape.left or 0) * sx)
    y = int((shape.top or 0) * sy)
    w = int((shape.width or 0) * sx)
    h = int((shape.height or 0) * sy)
    if w <= 0 or h <= 0:
        return

    # Handle groups recursively
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            try:
                _render_shape_to_image(draw, img, child, sx, sy)
            except Exception:
                continue
        return

    # Handle pictures
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            embedded = Image.open(BytesIO(shape.image.blob))
            embedded = embedded.resize((w, h), Image.LANCZOS)
            img.paste(embedded, (x, y))
        except Exception:
            draw.rectangle([x, y, x + w, y + h], fill=(230, 230, 230))
        return

    # Handle tables
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        _render_table_to_image(draw, shape.table, x, y, w, h)
        return

    # Shape fill
    fill = None
    try:
        if shape.fill and shape.fill.type is not None:
            fill = _hex_to_rgb_tuple(str(shape.fill.fore_color.rgb))
    except Exception:
        pass

    if fill:
        try:
            draw.rounded_rectangle([x, y, x + w, y + h], radius=min(6, w // 10, h // 10), fill=fill)
        except AttributeError:
            draw.rectangle([x, y, x + w, y + h], fill=fill)

    # Render text
    if shape.has_text_frame:
        ty = y + 4
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                text = run.text.strip()
                if not text:
                    continue
                fs = 14
                try:
                    if run.font.size:
                        fs = max(8, min(40, int(run.font.size.pt * (THUMB_H / 540))))
                except Exception:
                    pass
                fc = (50, 50, 50)
                try:
                    if run.font.color and run.font.color.rgb:
                        fc = _hex_to_rgb_tuple(str(run.font.color.rgb))
                except Exception:
                    pass
                bold = False
                try:
                    bold = run.font.bold or False
                except Exception:
                    pass

                font = _get_font(fs, bold)
                tx = x + 6
                if ty + fs < y + h:
                    draw.text((tx, ty), text[:60], fill=fc, font=font)
                ty += fs + 3
            ty += 3


def _render_table_to_image(draw, table, x, y, w, h):
    """Render a PPTX table onto the image."""
    rows = len(table.rows)
    cols = len(table.columns)
    if rows == 0 or cols == 0:
        return
    cw = w // cols
    ch = h // rows

    for ri, row in enumerate(table.rows):
        for ci, cell in enumerate(row.cells):
            cx = x + ci * cw
            cy = y + ri * ch
            fill = (255, 255, 255)
            try:
                if cell.fill and cell.fill.type is not None:
                    fill = _hex_to_rgb_tuple(str(cell.fill.fore_color.rgb))
                elif ri == 0:
                    fill = (0, 51, 141)
                elif ri % 2 == 1:
                    fill = (245, 245, 245)
            except Exception:
                if ri == 0:
                    fill = (0, 51, 141)
            draw.rectangle([cx, cy, cx + cw, cy + ch], fill=fill, outline=(200, 200, 200))
            text = cell.text.strip()
            if text:
                fs = max(7, min(12, ch // 3))
                fc = (255, 255, 255) if ri == 0 else (50, 50, 50)
                font = _get_font(fs)
                max_c = max(3, cw // (fs // 2 + 1))
                dt = text[:max_c] + "..." if len(text) > max_c else text
                draw.text((cx + 3, cy + (ch - fs) // 2), dt, fill=fc, font=font)


def _generate_placeholder(output_dir: str, index: int) -> str:
    """Generate a simple placeholder thumbnail."""
    path = os.path.join(output_dir, f"slide_{index}.png")
    img = Image.new("RGB", (THUMB_W, THUMB_H), (240, 242, 245))
    draw = ImageDraw.Draw(img)
    draw.rectangle([10, 10, THUMB_W - 10, THUMB_H - 10], outline=(200, 200, 200), width=2)
    font = _get_font(28)
    text = f"Slide {index + 1}"
    bbox = draw.textbbox((0, 0), text, font=font)
    tw = bbox[2] - bbox[0]
    draw.text(((THUMB_W - tw) // 2, THUMB_H // 2 - 14), text, fill=(150, 150, 150), font=font)
    img.save(path, "PNG")
    return path


# Keep old function name as alias for backward compatibility
def generate_thumbnail_from_design(design: dict, output_path: str, index: int) -> str:
    """Legacy: generate thumbnail from design dict. Now uses shape rendering."""
    img = Image.new("RGB", (THUMB_W, THUMB_H), (255, 255, 255))
    draw = ImageDraw.Draw(img)
    shapes = design.get("shapes", [])
    slide_w_in, slide_h_in = 13.333, 7.5

    for shape in shapes:
        pos = shape.get("position", {})
        left = pos.get("left_inches", 0)
        top = pos.get("top_inches", 0)
        w = pos.get("width_inches", 1)
        h = pos.get("height_inches", 1)
        x1 = max(0, int(left / slide_w_in * THUMB_W))
        y1 = max(0, int(top / slide_h_in * THUMB_H))
        x2 = min(THUMB_W, int((left + w) / slide_w_in * THUMB_W))
        y2 = min(THUMB_H, int((top + h) / slide_h_in * THUMB_H))
        if x2 <= x1 or y2 <= y1:
            continue

        if shape.get("fill_color"):
            fill = _hex_to_rgb_tuple(shape["fill_color"])
            try:
                draw.rounded_rectangle([x1, y1, x2, y2], radius=min(6, (x2 - x1) // 10), fill=fill)
            except AttributeError:
                draw.rectangle([x1, y1, x2, y2], fill=fill)

        text = shape.get("text", "")
        if text:
            fs = shape.get("font_size_pt") or 12
            scaled = max(8, min(32, int(fs * THUMB_W / (slide_w_in * 72))))
            font = _get_font(scaled, shape.get("font_bold", False))
            fc = _hex_to_rgb_tuple(shape.get("font_color")) if shape.get("font_color") else (50, 50, 50)
            draw.text((x1 + 4, y1 + 2), text[:50], fill=fc, font=font)

    if not shapes:
        font = _get_font(24)
        draw.text((THUMB_W // 2 - 60, THUMB_H // 2 - 12), f"Slide {index + 1}", fill=(180, 180, 180), font=font)

    draw.rectangle([0, 0, THUMB_W - 1, THUMB_H - 1], outline=(230, 230, 230), width=1)
    img.save(output_path, "PNG", quality=90)
    return output_path


def extract_slide_design(slide, prs) -> dict:
    """Extract design properties from a single PPTX slide."""
    design = {
        "background": {},
        "shapes": [],
        "color_palette": [],
        "content_slots": [],
        "estimated_items": 0,
        "supports_variable_items": False,
    }

    colors = set()
    total_shapes = len(list(slide.shapes))

    for shape_idx, shape in enumerate(slide.shapes):
        shape_info = {
            "type": "shape",
            "position": {
                "left_inches": _emu_to_inches(shape.left),
                "top_inches": _emu_to_inches(shape.top),
                "width_inches": _emu_to_inches(shape.width),
                "height_inches": _emu_to_inches(shape.height),
            },
            "fill_color": _extract_fill_color(shape),
            "rotation": shape.rotation if hasattr(shape, "rotation") else 0,
            "text": "",
            "font_size_pt": None,
            "font_bold": False,
            "font_color": None,
            "font_name": None,
            "alignment": None,
            "is_placeholder": False,
        }

        if shape.has_text_frame:
            tf = shape.text_frame
            full_text = tf.text.strip()
            shape_info["text"] = full_text
            if tf.paragraphs:
                para = tf.paragraphs[0]
                if para.runs:
                    run = para.runs[0]
                    shape_info["font_size_pt"] = run.font.size.pt if run.font.size else None
                    shape_info["font_bold"] = run.font.bold or False
                    shape_info["font_color"] = _color_to_hex(run.font.color) if run.font.color else None
                    shape_info["font_name"] = run.font.name
                shape_info["alignment"] = str(para.alignment) if para.alignment else None

            shape_info["is_placeholder"] = _is_placeholder_text(full_text)
            slot_type = _identify_slot_type(shape_info, shape_idx, total_shapes)
            if slot_type:
                design["content_slots"].append({
                    "slot_type": slot_type,
                    "shape_index": shape_idx,
                    "placeholder_text": full_text,
                })

        if shape_info["fill_color"]:
            colors.add(shape_info["fill_color"])
        if shape_info["font_color"]:
            colors.add(shape_info["font_color"])

        design["shapes"].append(shape_info)

    design["color_palette"] = list(colors)
    design["layout_style"] = _detect_layout_style(design["shapes"])
    item_slots = [s for s in design["content_slots"] if s["slot_type"] == "item"]
    design["estimated_items"] = len(item_slots)
    design["supports_variable_items"] = len(item_slots) > 2

    return design


async def process_template_upload(
    file_path: str,
    collection_name: str,
    description: str | None,
    icon: str | None,
    color: str | None,
    user_id: uuid.UUID,
    db: AsyncSession,
) -> dict:
    """Main entry: reads PPTX, extracts all slides as variations."""
    from app.services.pptx_parser import PPTXSlideParser

    prs = PptxPresentation(file_path)
    slide_count = len(prs.slides)

    collection = TemplateCollection(
        id=uuid.uuid4(),
        user_id=user_id,
        name=collection_name,
        description=description,
        icon=icon,
        color=color,
        source_filename=os.path.basename(file_path),
        source_file_path=file_path,
        variation_count=slide_count,
    )
    db.add(collection)

    # Deep parse all slides
    parser = PPTXSlideParser(file_path)
    parsed_slides = parser.parse_all_slides()
    logger.info("Parsed %d slides from %s", len(parsed_slides), file_path)

    # Generate ALL thumbnails (LibreOffice → PDF → PNG, or Python fallback)
    thumb_paths = generate_thumbnails_batch(file_path, str(collection.id), slide_count)

    variations = []
    for idx, slide in enumerate(prs.slides):
        design = extract_slide_design(slide, prs)
        slide_xml = etree.tostring(slide._element, pretty_print=True).decode()

        # Get parsed objects for this slide (strip shape_xml to keep JSON small)
        parsed = parsed_slides[idx] if idx < len(parsed_slides) else None
        objects_data = None
        if parsed:
            objects_data = {
                "slide_index": parsed["slide_index"],
                "slide_dimensions": parsed["slide_dimensions"],
                "background": parsed["background"],
                "layout_name": parsed.get("layout_name"),
                "objects": parsed["objects"],
                "object_count": parsed["object_count"],
                "color_palette": parsed["color_palette"],
                "font_inventory": parsed["font_inventory"],
                "has_images": parsed["has_images"],
                "has_charts": parsed["has_charts"],
                "has_tables": parsed["has_tables"],
            }

        # Use pre-generated thumbnail — convert full path to URL
        thumb_full = thumb_paths[idx] if idx < len(thumb_paths) else ""
        thumb_url = f"/uploads/templates/{collection.id}/thumbnails/slide_{idx}.png" if thumb_full else ""

        # Compute quality metrics from parsed objects + design slots
        metrics_input = dict(objects_data) if objects_data else {}
        if "content_slots" not in metrics_input and design.get("content_slots"):
            metrics_input["content_slots"] = design["content_slots"]
        metrics = compute_variation_metrics(metrics_input if metrics_input else None)

        variation = TemplateVariation(
            id=uuid.uuid4(),
            collection_id=collection.id,
            variation_index=idx,
            variation_name=_auto_name(design, idx),
            thumbnail_path=thumb_url,
            design_json=design,
            objects_json=objects_data,
            metrics_json=metrics,
            pptx_slide_xml=slide_xml,
            tags=_auto_tag(design),
        )
        db.add(variation)
        variations.append(variation)

    await db.flush()

    return {
        "collection_id": str(collection.id),
        "name": collection.name,
        "variation_count": len(variations),
    }

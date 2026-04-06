"""Deep PPTX slide parser — extracts complete object tree from every slide."""

import base64
import logging
from io import BytesIO

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Emu

logger = logging.getLogger("pptx_parser")


class PPTXSlideParser:
    """Parses every object on a PPTX slide into structured dicts."""

    def __init__(self, pptx_path: str):
        self.prs = Presentation(pptx_path)
        self.slide_w = self.prs.slide_width or Emu(12192000)
        self.slide_h = self.prs.slide_height or Emu(6858000)

    def parse_all_slides(self) -> list[dict]:
        slides = []
        for idx, slide in enumerate(self.prs.slides):
            slides.append(self.parse_slide(slide, idx))
        return slides

    def parse_slide(self, slide, index: int) -> dict:
        sd = {
            "slide_index": index,
            "slide_dimensions": {
                "width_inches": round(self.slide_w / 914400, 2),
                "height_inches": round(self.slide_h / 914400, 2),
            },
            "background": self._parse_bg(slide),
            "layout_name": slide.slide_layout.name if slide.slide_layout else None,
            "objects": [],
            "object_count": 0,
            "text_content": "",
            "color_palette": set(),
            "font_inventory": set(),
            "has_images": False,
            "has_charts": False,
            "has_tables": False,
        }
        texts = []
        for shape in slide.shapes:
            obj = self._parse_shape(shape, sd)
            if obj:
                sd["objects"].append(obj)
                if obj.get("text_content"):
                    texts.append(obj["text_content"])
        sd["object_count"] = len(sd["objects"])
        sd["text_content"] = "\n".join(texts)
        sd["color_palette"] = list(sd["color_palette"])
        sd["font_inventory"] = list(sd["font_inventory"])
        return sd

    def _parse_bg(self, slide) -> dict:
        bg = {"type": "none", "color": None}
        try:
            if slide.background and slide.background.fill and slide.background.fill.type is not None:
                ft = str(slide.background.fill.type)
                if "SOLID" in ft:
                    bg["type"] = "solid"
                    bg["color"] = f"#{slide.background.fill.fore_color.rgb}"
                elif "GRADIENT" in ft:
                    bg["type"] = "gradient"
        except Exception:
            pass
        return bg

    def _parse_shape(self, shape, sd: dict) -> dict | None:
        obj = {
            "id": shape.shape_id,
            "name": shape.name,
            "type": self._type_name(shape),
            "position": {
                "left": round((shape.left or 0) / 914400, 3),
                "top": round((shape.top or 0) / 914400, 3),
                "width": round((shape.width or 0) / 914400, 3),
                "height": round((shape.height or 0) / 914400, 3),
            },
            "rotation": getattr(shape, "rotation", 0),
            "fill": self._parse_fill(shape, sd),
            "border": self._parse_border(shape, sd),
            "text_content": None,
            "text_frames": None,
            "auto_shape": None,
            "image_data": None,
            "table_data": None,
            "chart_data": None,
            "group_children": None,
        }

        try:
            if hasattr(shape, "auto_shape_type") and shape.auto_shape_type:
                obj["auto_shape"] = str(shape.auto_shape_type)
        except Exception:
            pass

        # Text
        if shape.has_text_frame:
            td = self._parse_text_frame(shape.text_frame, sd)
            obj["text_frames"] = td["frames"]
            obj["text_content"] = td["text"]

        # Table
        st = getattr(shape, "shape_type", None)
        if st == MSO_SHAPE_TYPE.TABLE:
            sd["has_tables"] = True
            obj["table_data"] = self._parse_table(shape.table, sd)

        # Picture
        if st == MSO_SHAPE_TYPE.PICTURE:
            sd["has_images"] = True
            obj["image_data"] = self._parse_image(shape)

        # Chart
        if shape.has_chart:
            sd["has_charts"] = True
            obj["chart_data"] = self._parse_chart(shape.chart)

        # Group
        if st == MSO_SHAPE_TYPE.GROUP:
            children = []
            for child in shape.shapes:
                c = self._parse_shape(child, sd)
                if c:
                    children.append(c)
            obj["group_children"] = children

        return obj

    def _type_name(self, shape) -> str:
        m = {
            MSO_SHAPE_TYPE.AUTO_SHAPE: "auto_shape",
            MSO_SHAPE_TYPE.CHART: "chart",
            MSO_SHAPE_TYPE.FREEFORM: "freeform",
            MSO_SHAPE_TYPE.GROUP: "group",
            MSO_SHAPE_TYPE.LINE: "line",
            MSO_SHAPE_TYPE.PICTURE: "picture",
            MSO_SHAPE_TYPE.TABLE: "table",
            MSO_SHAPE_TYPE.TEXT_BOX: "text_box",
            MSO_SHAPE_TYPE.PLACEHOLDER: "placeholder",
        }
        try:
            return m.get(shape.shape_type, f"other")
        except Exception:
            return "unknown"

    def _parse_fill(self, shape, sd: dict) -> dict:
        f = {"type": "none", "color": None}
        try:
            if shape.fill and shape.fill.type is not None:
                ft = str(shape.fill.type)
                if "SOLID" in ft:
                    f["type"] = "solid"
                    rgb = str(shape.fill.fore_color.rgb)
                    f["color"] = f"#{rgb}"
                    sd["color_palette"].add(f"#{rgb}")
                elif "GRADIENT" in ft:
                    f["type"] = "gradient"
        except Exception:
            pass
        return f

    def _parse_border(self, shape, sd: dict) -> dict:
        b = {"has_border": False, "color": None, "width": None}
        try:
            if shape.line and shape.line.fill and shape.line.fill.type is not None:
                b["has_border"] = True
                b["color"] = f"#{shape.line.color.rgb}"
                sd["color_palette"].add(b["color"])
                b["width"] = shape.line.width.pt if shape.line.width else None
        except Exception:
            pass
        return b

    def _parse_text_frame(self, tf, sd: dict) -> dict:
        result = {"frames": [], "text": ""}
        texts = []
        paras = []
        for para in tf.paragraphs:
            pd = {
                "alignment": str(para.alignment) if para.alignment else None,
                "level": para.level,
                "text": para.text,
                "runs": [],
            }
            for run in para.runs:
                rd = {
                    "text": run.text,
                    "font": run.font.name,
                    "size": run.font.size.pt if run.font.size else None,
                    "bold": run.font.bold,
                    "italic": run.font.italic,
                    "color": None,
                }
                try:
                    if run.font.color and run.font.color.rgb:
                        rd["color"] = f"#{run.font.color.rgb}"
                        sd["color_palette"].add(rd["color"])
                except Exception:
                    pass
                if run.font.name:
                    sd["font_inventory"].add(run.font.name)
                pd["runs"].append(rd)
            paras.append(pd)
            if para.text.strip():
                texts.append(para.text)
        result["frames"] = paras
        result["text"] = "\n".join(texts)
        return result

    def _parse_table(self, table, sd: dict) -> dict:
        td = {"rows": len(table.rows), "columns": len(table.columns), "cells": []}
        for ri, row in enumerate(table.rows):
            row_cells = []
            for ci, cell in enumerate(row.cells):
                cd = {"row": ri, "col": ci, "text": cell.text, "fill": None}
                try:
                    if cell.fill and cell.fill.type is not None and "SOLID" in str(cell.fill.type):
                        cd["fill"] = f"#{cell.fill.fore_color.rgb}"
                        sd["color_palette"].add(cd["fill"])
                except Exception:
                    pass
                row_cells.append(cd)
            td["cells"].append(row_cells)
        return td

    def _parse_image(self, shape) -> dict:
        id_ = {"content_type": None, "width_px": None, "height_px": None, "size_bytes": 0}
        try:
            img = shape.image
            id_["content_type"] = img.content_type
            blob = img.blob
            id_["size_bytes"] = len(blob)
            from PIL import Image as PILImage
            pil = PILImage.open(BytesIO(blob))
            id_["width_px"] = pil.width
            id_["height_px"] = pil.height
            # Only store base64 for small images (<100KB)
            if len(blob) < 100000:
                id_["base64"] = base64.b64encode(blob).decode("utf-8")
        except Exception as e:
            logger.debug("Image parse error: %s", e)
        return id_

    def _parse_chart(self, chart) -> dict:
        cd = {"chart_type": str(chart.chart_type) if chart.chart_type else None, "series_count": 0}
        try:
            cd["series_count"] = len(list(chart.series))
        except Exception:
            pass
        try:
            cd["categories"] = [str(c) for c in chart.plots[0].categories]
        except Exception:
            pass
        return cd
